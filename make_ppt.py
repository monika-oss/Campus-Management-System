# ============================================================
#  Smart Campus Management System - PPT Generator
#  Run: python make_ppt.py
#  Screenshots folder: ppt_screenshots/
# ============================================================

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os, sys

# ── Paths ────────────────────────────────────────────────────
SS_DIR      = os.path.join(os.path.dirname(__file__), "ppt_screenshots")
OUTPUT_PATH = r"C:\Users\preet\OneDrive\Desktop\Smart_Campus_Presentation.pptx"

# ── Color Palette ─────────────────────────────────────────────
DARK_NAVY   = RGBColor(0x0F, 0x17, 0x2A)
MID_BLUE    = RGBColor(0x1E, 0x3A, 0x5F)
ACCENT_BLUE = RGBColor(0x25, 0x63, 0xEB)
LIGHT_BLUE  = RGBColor(0xDB, 0xEA, 0xFE)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
SLATE_400   = RGBColor(0x94, 0xA3, 0xB8)
SLATE_500   = RGBColor(0x64, 0x74, 0x8B)
SLATE_700   = RGBColor(0x33, 0x4E, 0x68)
SLATE_100   = RGBColor(0xF1, 0xF5, 0xF9)
GREEN       = RGBColor(0x05, 0x96, 0x69)
AMBER       = RGBColor(0xD9, 0x77, 0x06)
PURPLE      = RGBColor(0x7C, 0x3A, 0xED)
RED         = RGBColor(0xDC, 0x26, 0x26)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────
def rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = lw
    return s

def text(slide, txt, l, t, w, h, size=Pt(14), bold=False,
         color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = txt
    r.font.size = size
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return tb

def ss(name):
    """Return full path to a screenshot file, or None if missing."""
    p = os.path.join(SS_DIR, name)
    return p if os.path.exists(p) else None

def img_slide(title, img_file, slide_num, subtitle=""):
    """Slide with a full-bleed screenshot + title overlay bar."""
    s = prs.slides.add_slide(BLANK)
    # Background
    rect(s, 0, 0, 13.33, 7.5, fill=SLATE_100)
    # Image
    path = ss(img_file)
    if path:
        s.shapes.add_picture(path, Inches(0), Inches(1.0), Inches(13.33), Inches(6.5))
    else:
        # Placeholder when screenshot is missing
        rect(s, 0.3, 1.1, 12.73, 6.3, fill=MID_BLUE)
        text(s, f"[ Screenshot: {img_file} ]", 2, 3.5, 9.33, 1,
             size=Pt(18), bold=False, color=SLATE_400, align=PP_ALIGN.CENTER)
        text(s, f"Save your screenshot as:\nppt_screenshots/{img_file}",
             2, 4.5, 9.33, 1.2, size=Pt(14), bold=False,
             color=RGBColor(0x60,0x80,0xB0), align=PP_ALIGN.CENTER)
    # Title bar (on top of image)
    rect(s, 0, 0, 13.33, 1.0, fill=DARK_NAVY)
    text(s, title, 0.3, 0.08, 11.5, 0.85, size=Pt(28), bold=True,
         color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        text(s, subtitle, 0.3, 0.62, 9, 0.45, size=Pt(13),
             color=SLATE_400, align=PP_ALIGN.LEFT)
    text(s, str(slide_num), 12.8, 0.05, 0.5, 0.4,
         size=Pt(13), color=SLATE_500, align=PP_ALIGN.RIGHT)
    return s

def content_slide(num, title, subtitle, bullets, accent=ACCENT_BLUE, bg=DARK_NAVY):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.33, 7.5, fill=bg)
    rect(s, 0, 0, 13.33, 0.07, fill=accent)
    rect(s, 0, 0.07, 0.06, 7.43, fill=accent)
    text(s, title, 0.35, 0.12, 12.5, 0.9, size=Pt(34), bold=True, color=WHITE)
    if subtitle:
        text(s, subtitle, 0.35, 1.0, 12.2, 0.55, size=Pt(15), color=SLATE_400)
    y = 1.6
    for b in bullets:
        if isinstance(b, tuple):
            head, body = b
            text(s, f"  ->  {head}", 0.35, y, 12.5, 0.45,
                 size=Pt(17), bold=True, color=accent)
            y += 0.43
            text(s, f"        {body}", 0.35, y, 12.5, 0.44,
                 size=Pt(14), color=WHITE)
            y += 0.46
        else:
            text(s, f"   *   {b}", 0.35, y, 12.5, 0.48,
                 size=Pt(16), color=WHITE)
            y += 0.5
    text(s, str(num), 12.7, 7.1, 0.5, 0.35,
         size=Pt(12), color=SLATE_500, align=PP_ALIGN.RIGHT)
    return s

def three_col(num, title, cols, bg=DARK_NAVY):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.33, 7.5, fill=bg)
    rect(s, 0, 0, 13.33, 0.07, fill=ACCENT_BLUE)
    text(s, title, 0.35, 0.12, 12.5, 0.9, size=Pt(32), bold=True, color=WHITE)
    xstart, cw, gap = 0.4, 4.0, 0.26
    for i, (heading, col_color, items) in enumerate(cols):
        x = xstart + i*(cw+gap)
        rect(s, x, 1.15, cw, 5.9, fill=MID_BLUE, line=col_color, lw=Pt(1.5))
        rect(s, x, 1.15, cw, 0.65, fill=col_color)
        text(s, heading, x+0.1, 1.2, cw-0.2, 0.58, size=Pt(16), bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
        y = 2.0
        for item in items:
            text(s, f"  + {item}", x+0.15, y, cw-0.3, 0.55,
                 size=Pt(13), color=SLATE_400)
            y += 0.56
    text(s, str(num), 12.7, 7.1, 0.5, 0.35,
         size=Pt(12), color=SLATE_500, align=PP_ALIGN.RIGHT)
    return s


# ============================================================
#  SLIDE 1 – TITLE
# ============================================================
s1 = prs.slides.add_slide(BLANK)
rect(s1, 0, 0, 13.33, 7.5, fill=DARK_NAVY)
rect(s1, 0, 0, 13.33, 0.07, fill=ACCENT_BLUE)
rect(s1, 0, 7.43, 13.33, 0.07, fill=ACCENT_BLUE)
# Screenshot of login page as bg (optional)
path = ss("01_login.png")
if path:
    pic = s1.shapes.add_picture(path, Inches(7.0), Inches(1.2), Inches(5.9), Inches(5.0))
# Title text
text(s1, "Smart Campus", 0.5, 1.5, 7, 1.2, size=Pt(52), bold=True,
     color=WHITE, align=PP_ALIGN.LEFT)
text(s1, "Management System", 0.5, 2.7, 7, 1.1, size=Pt(40), bold=True,
     color=ACCENT_BLUE, align=PP_ALIGN.LEFT)
rect(s1, 0.5, 3.85, 3.5, 0.06, fill=ACCENT_BLUE)
text(s1, "Full Stack Web Development Project", 0.5, 4.0, 7, 0.65,
     size=Pt(18), color=SLATE_400, align=PP_ALIGN.LEFT)
text(s1, "Python Django  |  HTML/CSS/JS  |  Tailwind CSS  |  REST API",
     0.5, 4.7, 7, 0.55, size=Pt(14), color=RGBColor(0x60,0x80,0xB0))
text(s1, "Submitted to: Prime Vector\nTraining Program – Full Stack Development",
     0.5, 5.6, 7, 0.85, size=Pt(13), color=SLATE_500)
text(s1, "1", 12.7, 7.1, 0.5, 0.35, size=Pt(12), color=SLATE_500, align=PP_ALIGN.RIGHT)

# ============================================================
#  SLIDE 2 – AGENDA
# ============================================================
s2 = prs.slides.add_slide(BLANK)
rect(s2, 0, 0, 13.33, 7.5, fill=WHITE)
rect(s2, 0, 0, 0.5, 7.5, fill=ACCENT_BLUE)
text(s2, "Agenda", 0.8, 0.2, 11, 0.9, size=Pt(36), bold=True,
     color=DARK_NAVY, align=PP_ALIGN.LEFT)
items = [
    ("01", "Project Overview & Objectives"),
    ("02", "Problem Statement"),
    ("03", "Technology Stack"),
    ("04", "System Architecture"),
    ("05", "Role-Based Access Control"),
    ("06", "Admin Module – Dashboard, Students, Leave, Results"),
    ("07", "Faculty Module – Attendance, Leave Approvals"),
    ("08", "Student Module – Dashboard, Leaves, Results"),
    ("09", "Key Features & Security"),
    ("10", "Conclusion & Thank You"),
]
col = 0
for i, (num, label) in enumerate(items):
    x = 0.9 if col == 0 else 7.1
    y = 1.3 + (i % 5) * 1.05
    if i == 5: col = 1
    rect(s2, x, y, 0.55, 0.55, fill=ACCENT_BLUE)
    text(s2, num, x+0.02, y+0.05, 0.5, 0.45, size=Pt(16), bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
    text(s2, label, x+0.65, y+0.06, 5.5, 0.45, size=Pt(14),
         color=DARK_NAVY, align=PP_ALIGN.LEFT)
text(s2, "2", 12.7, 7.1, 0.5, 0.35, size=Pt(12), color=SLATE_500, align=PP_ALIGN.RIGHT)

# ============================================================
#  SLIDE 3 – PROBLEM STATEMENT
# ============================================================
content_slide(3, "Problem Statement",
    "What challenges exist in traditional campus management?",
    [
        ("Manual Attendance Tracking",
         "Faculty marks attendance on paper. Error-prone, time-consuming, no digital record."),
        ("No Digital Leave Workflow",
         "Students submit leave on paper. No multi-level approval. No tracking or history."),
        ("Scattered Academic Data",
         "Results, announcements, and attendance stored in separate unconnected systems."),
        ("No Role-Based Access",
         "No secure, separate portals for Admin / Faculty / Student with proper permissions."),
    ], accent=RED, bg=DARK_NAVY
)

# ============================================================
#  SLIDE 4 – PROJECT OVERVIEW
# ============================================================
content_slide(4, "Project Overview",
    "Smart Campus Management System – A centralized full-stack web platform",
    [
        "Single platform for Admins, Faculty, and Students",
        "Digitizes attendance, leave, results, and announcements",
        "Secure JWT-based login with 3 role-based portals",
        "Frontend hosted on Vercel | Backend REST API on Render cloud",
        "Responsive UI that works on Desktop, Tablet, and Mobile",
        "Real-time data from Django REST Framework backend",
    ]
)

# ============================================================
#  SLIDE 5 – TECHNOLOGY STACK
# ============================================================
s5 = prs.slides.add_slide(BLANK)
rect(s5, 0, 0, 13.33, 7.5, fill=DARK_NAVY)
rect(s5, 0, 0, 13.33, 0.07, fill=ACCENT_BLUE)
text(s5, "Technology Stack", 0.35, 0.12, 12.5, 0.9, size=Pt(34), bold=True, color=WHITE)
# Two-column grid of tech cards
categories = [
    ("Frontend", ACCENT_BLUE, ["HTML5 + CSS3", "JavaScript (ES6+)", "Tailwind CSS", "Lucide Icons"]),
    ("Backend",  GREEN,       ["Python 3.11", "Django 4.x", "Django REST Framework", "JWT Auth (SimpleJWT)"]),
    ("Database", PURPLE,      ["SQLite (dev)", "PostgreSQL (prod)", "Django ORM", "Migration system"]),
    ("Deploy",   AMBER,       ["Vercel (Frontend)", "Render.com (Backend)", "GitHub Pages", "GitHub CI/CD"]),
    ("Tools",    RED,         ["VS Code", "Git + GitHub", "Postman API testing", "Figma (wireframes)"]),
]
cols_per_row = 3
xstart, ystart, cw, ch, gap = 0.35, 1.2, 3.9, 2.6, 0.27
for i, (cat, color, items_list) in enumerate(categories):
    row, col = divmod(i, cols_per_row)
    x = xstart + col*(cw+gap)
    y = ystart + row*(ch+0.25)
    rect(s5, x, y, cw, ch, fill=MID_BLUE, line=color, lw=Pt(1.5))
    rect(s5, x, y, cw, 0.5, fill=color)
    text(s5, cat, x+0.1, y+0.05, cw-0.2, 0.42,
         size=Pt(15), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items_list):
        text(s5, f"  - {item}", x+0.12, y+0.55+j*0.47, cw-0.25, 0.45,
             size=Pt(12), color=SLATE_400)
text(s5, "5", 12.7, 7.1, 0.5, 0.35, size=Pt(12), color=SLATE_500, align=PP_ALIGN.RIGHT)

# ============================================================
#  SLIDE 6 – SYSTEM ARCHITECTURE
# ============================================================
s6 = prs.slides.add_slide(BLANK)
rect(s6, 0, 0, 13.33, 7.5, fill=DARK_NAVY)
rect(s6, 0, 0, 13.33, 0.07, fill=ACCENT_BLUE)
text(s6, "System Architecture", 0.35, 0.12, 12.5, 0.9, size=Pt(34), bold=True, color=WHITE)
# 3-Tier Diagram
tiers = [
    (ACCENT_BLUE, "TIER 1 – Frontend (Client)", "HTML / CSS / JavaScript / Tailwind", 1.1),
    (GREEN,       "TIER 2 – Backend API (Server)", "Python Django REST Framework + JWT Auth", 2.9),
    (PURPLE,      "TIER 3 – Database", "PostgreSQL / SQLite via Django ORM", 4.7),
]
for color, heading, sub, y in tiers:
    rect(s6, 2.5, y, 8.33, 1.45, fill=MID_BLUE, line=color, lw=Pt(2))
    rect(s6, 2.5, y, 8.33, 0.5, fill=color)
    text(s6, heading, 2.65, y+0.06, 8, 0.42, size=Pt(16), bold=True, color=WHITE)
    text(s6, sub, 2.65, y+0.55, 8, 0.55, size=Pt(13), color=SLATE_400)
# Arrows between tiers
for ay in [2.55, 4.37]:
    rect(s6, 6.3, ay, 0.73, 0.35, fill=ACCENT_BLUE)
    text(s6, " REST API", 6.32, ay+0.03, 0.8, 0.3, size=Pt(10), bold=True, color=WHITE)
# Users on left
text(s6, "USERS", 0.1, 2.0, 1.5, 0.5, size=Pt(13), bold=True, color=SLATE_400)
for uy, ul, uc in [(2.55, "Admin", ACCENT_BLUE), (3.3, "Faculty", GREEN), (4.05, "Student", AMBER)]:
    rect(s6, 0.1, uy, 1.4, 0.55, fill=MID_BLUE, line=uc, lw=Pt(1.5))
    text(s6, ul, 0.12, uy+0.06, 1.36, 0.42, size=Pt(13), bold=True, color=uc, align=PP_ALIGN.CENTER)
    # Arrow to frontend
    rect(s6, 1.5, uy+0.18, 1.0, 0.18, fill=uc)
# Cloud label
rect(s6, 10.5, 1.5, 2.4, 2.0, fill=MID_BLUE, line=AMBER, lw=Pt(1.5))
text(s6, "Cloud Deploy", 10.55, 1.55, 2.3, 0.45, size=Pt(12), bold=True, color=AMBER)
text(s6, "Vercel\n(Frontend)", 10.55, 2.0, 2.3, 0.6, size=Pt(11), color=SLATE_400)
text(s6, "Render.com\n(Backend)", 10.55, 2.65, 2.3, 0.6, size=Pt(11), color=SLATE_400)
text(s6, "6", 12.7, 7.1, 0.5, 0.35, size=Pt(12), color=SLATE_500, align=PP_ALIGN.RIGHT)

# ============================================================
#  SLIDE 7 – ROLE-BASED ACCESS CONTROL
# ============================================================
three_col(7, "Role-Based Access Control (JWT Authentication)",
    [
        ("Admin", ACCENT_BLUE, [
            "Manage Students & Faculty",
            "View all leave requests",
            "Upload exam results",
            "Post announcements",
            "Full system access",
        ]),
        ("Faculty", GREEN, [
            "Mark class attendance",
            "Approve / reject leaves",
            "Assign OD to students",
            "View own class results",
            "Department-limited view",
        ]),
        ("Student", AMBER, [
            "View own attendance",
            "Apply & track leaves",
            "View exam results",
            "Read announcements",
            "Personal portal only",
        ]),
    ]
)

# ============================================================
#  SLIDE 8 – ADMIN DASHBOARD (screenshot)
# ============================================================
img_slide("Admin Module – Dashboard", "02_admin_dashboard.png", 8,
          "KPI cards: Total Students, Faculty, Pending Leaves | Attendance Analytics")

# ============================================================
#  SLIDE 9 – ADMIN STUDENTS (screenshot)
# ============================================================
img_slide("Admin Module – Student Management", "03_admin_students.png", 9,
          "View, Add, Edit, Delete student records | Search & filter")

# ============================================================
#  SLIDE 10 – ADMIN LEAVE (screenshot)
# ============================================================
img_slide("Admin Module – Leave Management", "05_admin_leave.png", 10,
          "View all leave requests | Approve or Reject with status badges")

# ============================================================
#  SLIDE 11 – ADMIN RESULTS & ANNOUNCEMENTS
# ============================================================
three_col(11, "Admin Module – Results & Announcements",
    [
        ("Results Management", ACCENT_BLUE, [
            "Upload marks per subject",
            "Auto-calculate GPA/CGPA",
            "Grade letter assignment",
            "Semester-wise filtering",
            "Publish to student portal",
        ]),
        ("Announcements", GREEN, [
            "Create campus-wide notices",
            "Target specific departments",
            "Rich text formatting",
            "Timestamp & author info",
            "Real-time delivery",
        ]),
        ("Faculty Management", PURPLE, [
            "Add / edit faculty records",
            "Assign subjects to faculty",
            "Department grouping",
            "Role assignment & access",
            "View faculty activity",
        ]),
    ]
)

# ============================================================
#  SLIDE 12 – FACULTY DASHBOARD (screenshot)
# ============================================================
img_slide("Faculty Module – Dashboard", "08_faculty_dashboard.png", 12,
          "Quick stats: Today's Classes, Pending Approvals | Recent activity panel")

# ============================================================
#  SLIDE 13 – FACULTY ATTENDANCE (screenshot)
# ============================================================
img_slide("Faculty Module – Mark Attendance", "09_faculty_attendance.png", 13,
          "Select class + date -> Load Students -> Mark P/A/L/OD per student -> Submit")

# ============================================================
#  SLIDE 14 – FACULTY LEAVE APPROVALS (screenshot)
# ============================================================
img_slide("Faculty Module – Leave Approvals", "10_faculty_leave.png", 14,
          "Approve or Reject student leave requests | Multi-level workflow (Advisor -> HOD)")

# ============================================================
#  SLIDE 15 – FACULTY OD ASSIGNMENT
# ============================================================
content_slide(15, "Faculty Module – OD (On Duty) Assignment",
    "Assign On-Duty status to students for official activities",
    [
        ("Select Students", "Choose eligible students from your class list"),
        ("Enter OD Details", "Date, reason (seminar/workshop/sports/fest), duration"),
        ("Auto Attendance Update", "OD status auto-marked in attendance — NOT counted as absent"),
        ("Audit Trail", "Full history of OD assignments per faculty member"),
    ], accent=PURPLE
)

# ============================================================
#  SLIDE 16 – STUDENT DASHBOARD (screenshot)
# ============================================================
img_slide("Student Module – Dashboard", "11_student_dashboard.png", 16,
          "Welcome banner | Attendance % | Quick actions: Apply Leave, View Grades | Announcements")

# ============================================================
#  SLIDE 17 – STUDENT LEAVES (screenshot)
# ============================================================
img_slide("Student Module – Leave Applications", "12_student_leaves.png", 17,
          "Apply leave -> Advisor review -> HOD review -> Approved | Full history with status badges")

# ============================================================
#  SLIDE 18 – STUDENT RESULTS (screenshot)
# ============================================================
img_slide("Student Module – Academic Results", "13_student_results.png", 18,
          "Subject-wise marks | GPA / CGPA | Grade letters | Semester-wise view")

# ============================================================
#  SLIDE 19 – KEY FEATURES & SECURITY
# ============================================================
three_col(19, "Key Features, Security & Deployment",
    [
        ("Security & Auth", ACCENT_BLUE, [
            "JWT Token Authentication",
            "Role-based route protection",
            "CORS policy enforcement",
            "Secure password hashing",
            "Session auto-expiry",
        ]),
        ("UI/UX Highlights", GREEN, [
            "Responsive on all devices",
            "Glassmorphism design",
            "Dark sidebar navigation",
            "Toast notifications",
            "Accessible mobile UI",
        ]),
        ("Deployment Stack", AMBER, [
            "Frontend -> Vercel",
            "Backend -> Render.com",
            "DB -> PostgreSQL (cloud)",
            "Source -> GitHub",
            "Zero-downtime updates",
        ]),
    ]
)

# ============================================================
#  SLIDE 20 – THANK YOU
# ============================================================
s20 = prs.slides.add_slide(BLANK)
rect(s20, 0, 0, 13.33, 7.5, fill=DARK_NAVY)
rect(s20, 0, 0, 13.33, 0.07, fill=ACCENT_BLUE)
rect(s20, 0, 7.43, 13.33, 0.07, fill=ACCENT_BLUE)
# Screenshot collage on right (if available)
for idx, (fname, px, py, pw, ph) in enumerate([
    ("09_faculty_attendance.png", 7.0, 0.3, 6.0, 3.3),
    ("11_student_dashboard.png",  7.0, 3.7, 6.0, 3.5),
]):
    p = ss(fname)
    if p:
        s20.shapes.add_picture(p, Inches(px), Inches(py), Inches(pw), Inches(ph))
# Text
text(s20, "Thank You!", 0.5, 1.0, 6.3, 1.5,
     size=Pt(56), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
rect(s20, 0.5, 2.6, 3.0, 0.07, fill=ACCENT_BLUE)
text(s20, "Smart Campus Management System", 0.5, 2.75, 6.3, 0.7,
     size=Pt(20), bold=True, color=LIGHT_BLUE)
text(s20, "Full Stack Web Development Project\nSubmitted to: Prime Vector",
     0.5, 3.5, 6.3, 0.8, size=Pt(15), color=SLATE_400)
text(s20, "Tech Stack:", 0.5, 4.45, 2, 0.4, size=Pt(13), bold=True, color=SLATE_500)
text(s20, "Python | Django REST | HTML | CSS | JavaScript | Tailwind CSS",
     0.5, 4.85, 6.3, 0.5, size=Pt(13), color=SLATE_500)
rect(s20, 0.5, 5.55, 6.0, 0.75, fill=ACCENT_BLUE)
text(s20, "GitHub: github.com/monika-oss/Campus-Management-System",
     0.6, 5.62, 5.8, 0.6, size=Pt(13), bold=True, color=WHITE)
text(s20, "20", 12.7, 7.1, 0.5, 0.35, size=Pt(12), color=SLATE_500, align=PP_ALIGN.RIGHT)

# ============================================================
#  SAVE
# ============================================================
prs.save(OUTPUT_PATH)
print(f"PPT saved to: {OUTPUT_PATH}")
print(f"Total slides: {len(prs.slides)}")

# Check which screenshots are missing
missing = []
needed = [
    "01_login.png","02_admin_dashboard.png","03_admin_students.png",
    "04_admin_faculty.png","05_admin_leave.png","06_admin_results.png",
    "07_admin_announce.png","08_faculty_dashboard.png","09_faculty_attendance.png",
    "10_faculty_leave.png","11_student_dashboard.png","12_student_leaves.png",
    "13_student_results.png"
]
for n in needed:
    if not os.path.exists(os.path.join(SS_DIR, n)):
        missing.append(n)
if missing:
    print("\nMISSING screenshots (add to ppt_screenshots/ folder):")
    for m in missing:
        print(f"  - {m}")
else:
    print("\nAll screenshots found! PPT is complete.")
